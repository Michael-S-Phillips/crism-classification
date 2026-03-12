"""
Generate CRISM Classification summary PowerPoint presentation.

Produces a 7-slide deck (widescreen 16:9):
  1. Title
  2. Dataset Overview
  3. Methodology Journey (mrrsu → mrral, model progression)
  4. Representative Class Spectra (ratio to 'other')
  5. Hellas Domain Shift
  6. Ablation Results
  7. Key Findings & Next Steps

Usage:
    python scripts/generate_pptx.py
    python scripts/generate_pptx.py --out reports/my_deck.pptx
"""
import argparse
import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

PROJ = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
REPORTS = os.path.join(PROJ, 'reports')

# ── Colour palette ────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0d, 0x1b, 0x2a)
BLUE   = RGBColor(0x1f, 0x77, 0xb4)
ACCENT = RGBColor(0x4e, 0x9a, 0xf1)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGREY  = RGBColor(0xCC, 0xCC, 0xCC)
GOLD   = RGBColor(0xF4, 0xA7, 0x42)

# ── Slide dimensions (16:9) ───────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(layout)


def fill_bg(slide, color: RGBColor = NAVY):
    from pptx.oxml.ns import qn
    import lxml.etree as etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text: str, left, top, width, height,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def add_picture_centered(slide, img_path: str, top, height, left_margin=Inches(0.3)):
    """Add image maintaining aspect ratio, centred horizontally."""
    from PIL import Image
    with Image.open(img_path) as im:
        iw, ih = im.size
    aspect = iw / ih
    w = height * aspect
    left = (W - w) / 2
    if left < left_margin:
        w    = W - 2 * left_margin
        left = left_margin
        height = w / aspect
    slide.shapes.add_picture(img_path, left, top, width=w, height=height)


def accent_line(slide, left, top, width, color=ACCENT, height=Pt(2)):
    """Thin horizontal accent line."""
    from pptx.util import Emu
    ln = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, Emu(int(height))
    )
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()


# ── Slide builders ────────────────────────────────────────────────────────────

def slide_title(prs):
    sl = blank_slide(prs)
    fill_bg(sl)

    # Big title
    add_textbox(sl, 'CRISM MRDR Mineral Classification',
                Inches(1.2), Inches(1.8), Inches(10.9), Inches(1.4),
                font_size=36, bold=True, align=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(sl, 'Spectral Transformer + MAE Pre-training for Mars Mineral Mapping',
                Inches(1.2), Inches(3.0), Inches(10.9), Inches(0.7),
                font_size=18, color=LGREY, align=PP_ALIGN.CENTER, italic=True)

    accent_line(sl, Inches(2.0), Inches(3.9), Inches(9.33))

    # Stat block
    stats = [
        ('1.97 M', 'labeled pixels'),
        ('2 basins', 'Argyre + Hellas'),
        ('59 bands', 'mrral 410–2457 nm'),
        ('5 classes', 'Olivine · LCP · HCP · Plag · Other'),
    ]
    col_w = Inches(3.0)
    for i, (val, lab) in enumerate(stats):
        x = Inches(0.4) + i * col_w + (W - 4*col_w - Inches(0.8)) / 2
        add_textbox(sl, val,  x, Inches(4.2), col_w, Inches(0.6),
                    font_size=26, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_textbox(sl, lab,  x, Inches(4.8), col_w, Inches(0.5),
                    font_size=12, color=LGREY, align=PP_ALIGN.CENTER)

    add_textbox(sl, 'March 2026  |  Space Imagery Center',
                Inches(0), Inches(6.9), W, Inches(0.4),
                font_size=10, color=LGREY, align=PP_ALIGN.CENTER, italic=True)


def slide_dataset(prs):
    sl = blank_slide(prs)
    fill_bg(sl)

    add_textbox(sl, 'Dataset Overview', Inches(0.4), Inches(0.15), Inches(9), Inches(0.6),
                font_size=24, bold=True)
    accent_line(sl, Inches(0.4), Inches(0.72), Inches(12.5))

    img = os.path.join(REPORTS, 'fig_dataset_overview.png')
    if os.path.exists(img):
        add_picture_centered(sl, img, top=Inches(0.85), height=Inches(6.3))

    add_textbox(sl,
                'Argyre basin: hard labels (1.0/0.0)  ·  '
                'Hellas basin: soft labels (0.5) — olivine uncertainty from mixed spectral units  ·  '
                'Argyre HCP is predominantly co-labeled with olivine (Olivine+HCP)',
                Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.35),
                font_size=9, color=LGREY, italic=True)


def slide_methodology(prs):
    """Two-phase methodology journey: mrrsu param cubes → mrral full spectra."""
    sl = blank_slide(prs)
    fill_bg(sl)

    add_textbox(sl, 'Methodology Journey', Inches(0.4), Inches(0.15), Inches(9), Inches(0.6),
                font_size=24, bold=True)
    accent_line(sl, Inches(0.4), Inches(0.72), Inches(12.5))

    # Phase 1 box
    from pptx.util import Emu
    box1 = sl.shapes.add_shape(1, Inches(0.4), Inches(0.95), Inches(5.9), Inches(5.9))
    box1.fill.solid(); box1.fill.fore_color.rgb = RGBColor(0x10, 0x28, 0x40)
    box1.line.color.rgb = ACCENT; box1.line.width = Pt(1.2)

    add_textbox(sl, 'Phase 1 — mrrsu Parameter Cubes',
                Inches(0.55), Inches(1.0), Inches(5.6), Inches(0.45),
                font_size=14, bold=True, color=ACCENT)
    add_textbox(sl, '60 pre-computed band ratio / summary parameters\n(BD1300, OLINDEX3, LCPINDEX2, HCPINDEX2, …)',
                Inches(0.55), Inches(1.45), Inches(5.6), Inches(0.6),
                font_size=10, color=LGREY, italic=True)

    phase1_models = [
        ('Logistic Regression',         '0.560'),
        ('Random Forest',               '0.608'),
        ('XGBoost',                     '0.609'),
        ('LightGBM',                    '0.616'),
        ('MLP (pixel)',                 '0.648'),
        ('Spatial CNN  (7×7 patches)',  '0.672'),
        ('Spatial ViT  (7×7 patches)',  '0.675  ★ best'),
    ]
    for i, (name, score) in enumerate(phase1_models):
        y = Inches(2.1) + i * Inches(0.52)
        color = GOLD if '★' in score else WHITE
        add_textbox(sl, f'  {name}', Inches(0.55), y, Inches(4.0), Inches(0.45),
                    font_size=11, color=color)
        add_textbox(sl, f'mAP {score}', Inches(4.55), y, Inches(1.5), Inches(0.45),
                    font_size=11, bold='★' in score, color=color, align=PP_ALIGN.RIGHT)

    add_textbox(sl,
                'Limitation: band-ratio cubes lose spectral shape;\n'
                'absorption depth & width are pre-compressed.',
                Inches(0.55), Inches(5.7), Inches(5.6), Inches(0.8),
                font_size=9, color=LGREY, italic=True)

    # Arrow
    arrow_x = Inches(6.5)
    add_textbox(sl, '→', arrow_x, Inches(3.6), Inches(0.5), Inches(0.6),
                font_size=32, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(sl, 'Switch to\nfull spectra',
                arrow_x - Inches(0.05), Inches(4.2), Inches(0.65), Inches(0.7),
                font_size=9, color=LGREY, align=PP_ALIGN.CENTER, italic=True)

    # Phase 2 box
    box2 = sl.shapes.add_shape(1, Inches(7.1), Inches(0.95), Inches(5.85), Inches(5.9))
    box2.fill.solid(); box2.fill.fore_color.rgb = RGBColor(0x10, 0x28, 0x40)
    box2.line.color.rgb = GOLD; box2.line.width = Pt(1.2)

    add_textbox(sl, 'Phase 2 — mrral Full Reflectance Spectra',
                Inches(7.25), Inches(1.0), Inches(5.55), Inches(0.45),
                font_size=14, bold=True, color=GOLD)
    add_textbox(sl, '59 bands · 410–2457 nm · full reflectance (no pre-processing)',
                Inches(7.25), Inches(1.45), Inches(5.55), Inches(0.4),
                font_size=10, color=LGREY, italic=True)

    phase2_models = [
        ('SpectralCNN 1D  (pixel)',            '0.626'),
        ('SpectralCNN 1D  + focal loss',       '0.636'),
        ('SpectralViT  (no pretrain)',          '0.613'),
        ('SpectralViT  + MAE pretrain',        '0.621  ↑ plag +14 pp'),
        ('sweep_v5 (5-class, in progress)',    '…'),
    ]
    for i, (name, score) in enumerate(phase2_models):
        y = Inches(2.1) + i * Inches(0.65)
        color = GOLD if '↑' in score else WHITE
        add_textbox(sl, f'  {name}', Inches(7.25), y, Inches(4.3), Inches(0.55),
                    font_size=11, color=color)
        add_textbox(sl, f'mAP {score}', Inches(11.4), y, Inches(1.35), Inches(0.55),
                    font_size=11, bold='★' in score, color=color, align=PP_ALIGN.RIGHT)

    add_textbox(sl,
                'MAE pre-training recovers fine spectral absorptions.\n'
                'Focal loss helps with rare classes (HCP, plagioclase).',
                Inches(7.25), Inches(5.55), Inches(5.55), Inches(0.8),
                font_size=9, color=LGREY, italic=True)

    # Bottom note
    add_textbox(sl,
                'mrrsu spatial CNN/ViT used 7×7 pixel patch context; mrral models are per-pixel spectral (no spatial context yet)',
                Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.35),
                font_size=9, color=LGREY, italic=True)


def slide_spectra(prs):
    sl = blank_slide(prs)
    fill_bg(sl)

    add_textbox(sl, 'Representative Class Spectra  (ratio to "other" class)',
                Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.6),
                font_size=24, bold=True)
    accent_line(sl, Inches(0.4), Inches(0.72), Inches(12.5))

    img = os.path.join(REPORTS, 'class_spectra_ratio.png')
    if os.path.exists(img):
        add_picture_centered(sl, img, top=Inches(0.82), height=Inches(6.3))

    add_textbox(sl,
                'Median ± 10th–90th percentile  ·  '
                'Denominator: median of "other" class pixels (physically neutral background)  ·  '
                'High + Moderate confidence pixels only',
                Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.35),
                font_size=9, color=LGREY, italic=True)


def slide_domain_shift(prs):
    sl = blank_slide(prs)
    fill_bg(sl)

    add_textbox(sl, 'Hellas Domain Shift — Challenge & Findings',
                Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.6),
                font_size=24, bold=True)
    accent_line(sl, Inches(0.4), Inches(0.72), Inches(12.5))

    img = os.path.join(REPORTS, 'hellas_domain_shift.png')
    if os.path.exists(img):
        # leave room for bullet points on right
        sl.shapes.add_picture(img, Inches(0.3), Inches(0.85),
                               width=Inches(7.8), height=Inches(6.0))

    bullets = [
        ('Domain shift', 'Hellas pixels ~27% brighter (mean IF 0.194 vs 0.153 for Argyre)'),
        ('Balanced sampler ✗',
         'HCP is only 1.7% of Hellas pixels vs 21.6% for Argyre → combined dataset '
         'makes HCP very rare → 150× upweight → LCP/plagioclase AP collapses to 0.09'),
        ('Focal loss ✓',
         'Does not reweight by global class frequency; '
         'focuses on hard examples regardless of domain'),
        ('Argyre HCP purity',
         '94.7% of Argyre HCP pixels are co-labeled with olivine (Olivine+HCP); '
         'only 0.9% are pure HCP'),
        ('Hellas HCP quality',
         '100% of Hellas HCP pixels are pure — '
         'the only clean HCP training signal in the dataset'),
        ('5-class schema',
         'Collapse olivine_t1/t2 → olivine; uniform confidence weights (1.0); '
         'n_classes 6 → 5'),
    ]
    y = Inches(0.9)
    for title, body in bullets:
        add_textbox(sl, title, Inches(8.3), y, Inches(4.8), Inches(0.35),
                    font_size=11, bold=True, color=ACCENT)
        add_textbox(sl, body,  Inches(8.3), y + Inches(0.32), Inches(4.8), Inches(0.65),
                    font_size=9.5, color=LGREY)
        y += Inches(1.02)


def slide_ablation(prs):
    sl = blank_slide(prs)
    fill_bg(sl)

    add_textbox(sl, 'Ablation Results — mrral Spectral Models (sweep_v3)',
                Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.6),
                font_size=24, bold=True)
    accent_line(sl, Inches(0.4), Inches(0.72), Inches(12.5))

    img = os.path.join(REPORTS, 'fig_ablation.png')
    if os.path.exists(img):
        add_picture_centered(sl, img, top=Inches(0.85), height=Inches(6.1))

    add_textbox(sl,
                'Note: first 3 configs (scnn_aug, scnn_balanced, scnn_base) trained on 1.25M-pixel parquet; '
                'last 6 trained on 1.97M-pixel parquet (Hellas appended mid-sweep). '
                'sweep_v5 (5-class, focal loss, uniform weights) now running for clean comparison.',
                Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.4),
                font_size=8.5, color=LGREY, italic=True)


def slide_findings(prs):
    sl = blank_slide(prs)
    fill_bg(sl)

    add_textbox(sl, 'Key Findings & Next Steps',
                Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.6),
                font_size=24, bold=True)
    accent_line(sl, Inches(0.4), Inches(0.72), Inches(12.5))

    findings = [
        ('Full spectra > band ratios',
         'mrral 59-band reflectance lets models learn diagnostic absorption features '
         '(1 µm olivine band, 2 µm pyroxene band) that mrrsu band ratios compress away.'),
        ('MAE pre-training helps rare classes',
         'SpectralViT + MAE pretrain: plagioclase AP 0.419 → 0.562 (+14 pp). '
         'Unsupervised spectral reconstruction forces encoder to learn fine absorption geometry.'),
        ('Hellas HCP is the only clean HCP signal',
         'Argyre HCP is 94.7% co-labeled with olivine. The 256k Hellas HCP pixels are '
         'pure and critical for learning the 1 µm + 2.3 µm pyroxene absorptions.'),
        ('Balanced sampling fails with domain shift',
         '~27% brightness offset between Argyre and Hellas; class-frequency upweighting '
         'conflates domain imbalance with class imbalance. Focal loss is robust.'),
        ('5-class schema reduces label noise',
         'Collapsing olivine types and uniform confidence weights (1.0) give cleaner '
         'gradient signal; type discrimination deferred to a second-stage model.'),
    ]

    next_steps = [
        'Evaluate sweep_v5 (5-class, focal loss): compare scnn_base_v5 / svit_base_v5 / svit_mae_v5',
        'Ensemble top-3 checkpoints on locked test set for final numbers',
        'Add spatial context: 7×7 patch CNN/ViT operating on mrral spectra (vs mrrsu)',
        'Monte Carlo dropout or conformal prediction for per-pixel confidence calibration',
        'Olivine type 1 vs type 2 discrimination as second-stage model',
    ]

    # Left column: findings
    y = Inches(0.85)
    for title, body in findings:
        add_textbox(sl, f'● {title}', Inches(0.4), y, Inches(7.3), Inches(0.38),
                    font_size=11, bold=True, color=ACCENT)
        add_textbox(sl, body, Inches(0.6), y + Inches(0.36), Inches(7.1), Inches(0.55),
                    font_size=9.5, color=WHITE)
        y += Inches(1.0)

    # Right column: next steps
    accent_line(sl, Inches(8.0), Inches(0.85), Inches(0.0), color=ACCENT)  # placeholder
    add_textbox(sl, 'Next Steps', Inches(8.1), Inches(0.85), Inches(4.9), Inches(0.45),
                font_size=14, bold=True, color=GOLD)
    y2 = Inches(1.35)
    for step in next_steps:
        add_textbox(sl, f'→  {step}', Inches(8.1), y2, Inches(4.9), Inches(0.75),
                    font_size=10, color=LGREY)
        y2 += Inches(1.1)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument('--out', default=os.path.join(REPORTS, 'crism_classification_summary_v3.pptx'))
    args = parser.parse_args()

    prs = new_prs()
    print('Building slides...')
    slide_title(prs)       ; print('  1/7  Title')
    slide_dataset(prs)     ; print('  2/7  Dataset Overview')
    slide_methodology(prs) ; print('  3/7  Methodology Journey')
    slide_spectra(prs)     ; print('  4/7  Class Spectra')
    slide_domain_shift(prs); print('  5/7  Domain Shift')
    slide_ablation(prs)    ; print('  6/7  Ablation Results')
    slide_findings(prs)    ; print('  7/7  Key Findings & Next Steps')

    os.makedirs(os.path.dirname(args.out), exist_ok=True)
    prs.save(args.out)
    print(f'\nSaved → {args.out}')


if __name__ == '__main__':
    main()
